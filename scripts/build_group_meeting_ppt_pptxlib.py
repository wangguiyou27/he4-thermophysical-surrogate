from __future__ import annotations

import shutil
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


WIDE = (13.333333, 7.5)


def add_title(slide, title, subtitle=None):
    box = slide.shapes.add_textbox(Inches(0.45), Inches(0.18), Inches(12.45), Inches(0.55))
    p = box.text_frame.paragraphs[0]
    p.text = title
    p.font.name = "Microsoft YaHei"
    p.font.size = Pt(25)
    p.font.bold = True
    p.font.color.rgb = RGBColor(17, 24, 39)
    line = slide.shapes.add_shape(1, Inches(0.45), Inches(0.85), Inches(12.45), Inches(0.03))
    line.fill.solid()
    line.fill.fore_color.rgb = RGBColor(37, 99, 235)
    line.line.fill.background()
    if subtitle:
        sub = slide.shapes.add_textbox(Inches(0.48), Inches(0.9), Inches(12.2), Inches(0.35))
        p = sub.text_frame.paragraphs[0]
        p.text = subtitle
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(13)
        p.font.color.rgb = RGBColor(100, 116, 139)


def add_text(slide, x, y, w, h, text, size=22, color=(31, 41, 55), bold=False):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = RGBColor(*color)
        p.space_after = Pt(4)
    return box


def add_box(slide, x, y, w, h, text, size=21, fill=(238, 242, 255), color=(17, 24, 39), bold=False):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(*fill)
    shape.line.color.rgb = RGBColor(203, 213, 225)
    tf = shape.text_frame
    tf.clear()
    tf.margin_left = Inches(0.12)
    tf.margin_right = Inches(0.12)
    tf.margin_top = Inches(0.08)
    tf.margin_bottom = Inches(0.08)
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = RGBColor(*color)
        p.alignment = PP_ALIGN.CENTER
    return shape


def add_image(slide, path, x, y, max_w, max_h):
    path = Path(path)
    if not path.exists():
        add_text(slide, x, y, max_w, max_h, f"Missing image:\n{path}", size=14, color=(185, 28, 28))
        return
    slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(max_w), height=Inches(max_h))


def blank_slide(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, title, subtitle)
    return slide


def build(root: Path):
    prs = Presentation()
    prs.slide_width = Inches(WIDE[0])
    prs.slide_height = Inches(WIDE[1])
    fig = root / "paper_outputs" / "final_he4_single_phase" / "figures"

    s = blank_slide(prs, "面向低温工程计算的氦4单相热物性神经网络代理模型", "组会汇报 | REFPROP 数据集、多物性 ANN、速度与基线对比")
    add_text(s, 0.7, 1.45, 7.2, 1.4, "研究目标：建立一个可替代 REFPROP 高频调用的高精度、快速、可部署物性代理模型", 24, bold=True)
    add_box(s, 0.7, 3.1, 3.2, 1.0, "Fluid\nHelium-4", 22, (219, 234, 254), bold=True)
    add_box(s, 4.2, 3.1, 3.2, 1.0, "Domain\n5.5-300 K\n0.001-3 MPa", 20, (220, 252, 231), bold=True)
    add_box(s, 7.7, 3.1, 3.2, 1.0, "Outputs\n8 properties", 22, (254, 243, 199), bold=True)
    add_text(s, 0.7, 5.7, 9.0, 0.5, "当前状态：模型训练、误差分析、REFPROP 速度对比、传统 ML 基线对比已完成", 20, (71, 85, 105))

    s = blank_slide(prs, "研究背景与问题", "为什么需要 ANN 物性代理模型")
    add_text(s, 0.7, 1.35, 5.5, 4.7, "• 低温系统仿真中需要反复调用氦物性\n• REFPROP 精度高，但逐点调用在大规模迭代/优化中较慢\n• 单个坏点可能导致工程迭代发散\n• 目标：在保持 REFPROP 级精度的同时，提高调用速度，并分析误差风险区域", 22)
    add_box(s, 7.0, 1.55, 2.1, 0.85, "T, P", 26, (224, 242, 254), bold=True)
    add_box(s, 7.0, 2.8, 2.1, 0.85, "REFPROP", 24, (253, 230, 138), bold=True)
    add_box(s, 7.0, 4.05, 2.1, 0.85, "ANN surrogate", 21, (220, 252, 231), bold=True)
    add_box(s, 9.7, 2.8, 2.4, 2.1, "ρ, s, h, w\nCv, Cp, μ, k", 22, (238, 242, 255), bold=True)

    s = blank_slide(prs, "数据集与适用范围", "本文主线：He-4 单相、5.5 K 以上")
    add_text(s, 0.7, 1.28, 5.9, 4.9, "• 数据来源：NIST REFPROP 10\n• 相区：单相区\n• 温度范围：5.5-300 K\n• 压力范围：0.001-3 MPa\n• 样本数：296646\n• 输入：Temperature, Pressure\n• 输出：8 个热物性", 23)
    add_box(s, 7.0, 1.35, 4.8, 4.8, "输出物性\n\nDensity\nEntropy\nEnthalpy\nSpeed of sound\nCv / Cp\nViscosity\nThermal conductivity", 22, (248, 250, 252))

    s = blank_slide(prs, "模型架构", "多输出 MLP，同时预测 8 个物性")
    add_box(s, 0.7, 1.25, 2.0, 0.8, "T, P", 26, (224, 242, 254), bold=True)
    add_box(s, 3.1, 1.25, 2.5, 0.8, "log10 + scaler", 22, (219, 234, 254), bold=True)
    add_box(s, 6.0, 1.05, 3.0, 1.2, "MLP\n[256,512,512,256,128]", 20, (220, 252, 231), bold=True)
    add_box(s, 9.5, 1.25, 2.7, 0.8, "8 properties", 22, (254, 243, 199), bold=True)
    add_text(s, 0.9, 3.0, 11.2, 3.1, "网络模块：Linear → LayerNorm → GELU → Dropout\n优化器：AdamW，weight decay = 1e-5，cosine annealing\n训练策略：early stopping，多物性联合损失\n优势：结构紧凑，推理速度快，便于工程部署", 23)

    s = blank_slide(prs, "关键方法：目标变换与加权训练", "解决量纲差异、正值约束和低温困难区域")
    add_text(s, 0.7, 1.3, 6.0, 4.8, "目标变换\n• 正值物性：log10(y)\n• 熵：log10(S + 1.0)\n\n加权训练\n• 低温区 5.5-20 K：权重 4.0\n• 过渡区 20-80 K：权重 1.5\n• Density 权重 2.0\n• Entropy 权重 3.0", 22)
    add_text(s, 7.0, 1.4, 5.2, 4.6, "动机\n• 避免密度等正值物性出现非物理解\n• 缓解不同物性尺度差异\n• 强化低温区和困难目标\n• 降低尾部坏点对工程迭代的风险", 22, (51, 65, 85))

    s = blank_slide(prs, "最终模型总体精度", "全测试集统计")
    cards = [("MAPE", "0.150%"), ("Median", "0.072%"), ("P95", "0.541%"), ("P99", "1.126%"), ("Mean R²", "0.999951"), ("Within 5%", "99.998%")]
    fills = [(219, 234, 254), (220, 252, 231), (254, 243, 199), (237, 233, 254), (255, 228, 230), (236, 252, 203)]
    for i, (k, v) in enumerate(cards):
        add_box(s, 0.8 + (i % 3) * 4.1, 1.4 + (i // 3) * 2.0, 3.45, 1.25, f"{k}\n{v}", 25, fills[i], bold=True)
    add_text(s, 0.9, 5.8, 10.8, 0.6, "结论：整体误差低，P99 约 1.13%，绝大多数点落在 5% 以内。", 22)

    s = blank_slide(prs, "逐物性误差分析", "不同物性的预测难度不同")
    add_image(s, fig / "accuracy" / "property_error_bars.png", 0.6, 1.2, 6.1, 5.5)
    add_image(s, fig / "accuracy" / "property_accuracy_r2.png", 6.8, 1.2, 5.8, 5.5)

    s = blank_slide(prs, "温压平面误差分布与坏点分析", "定位模型风险区域")
    add_image(s, fig / "accuracy" / "error_maps" / "overall_error_heatmap.png", 0.6, 1.2, 5.8, 5.4)
    add_image(s, fig / "accuracy" / "error_maps" / "worst_points_tp_scatter.png", 6.7, 1.2, 5.8, 5.4)

    s = blank_slide(prs, "REFPROP vs ANN 速度对比", "大规模状态点调用")
    add_image(s, fig / "speed" / "refprop_ann_speed.png", 0.6, 1.2, 5.8, 5.4)
    add_image(s, fig / "speed" / "ann_speedup_over_refprop.png", 6.7, 1.2, 5.8, 5.4)
    add_text(s, 0.8, 6.75, 11.0, 0.35, "全数据 296646 点：ANN ≈ 418038 points/s，REFPROP ≈ 12113 points/s，加速约 34.5×", 16, (71, 85, 105))

    s = blank_slide(prs, "传统机器学习基线对比", "精度、速度与尾部误差的综合比较")
    add_image(s, fig / "baselines" / "baseline_accuracy_speed.png", 0.6, 1.15, 5.9, 5.45)
    add_image(s, fig / "baselines" / "baseline_tail_errors.png", 6.75, 1.15, 5.7, 5.45)

    s = blank_slide(prs, "基线对比结论", "不能只看平均误差")
    add_text(s, 0.75, 1.3, 11.7, 4.8, "• Extra Trees 的平均 MAPE 最低，但推理速度明显低于当前 MLP\n• Random Forest / Gaussian RBF Ridge 平均误差也较低，但最大坏点更大\n• 当前 MLP 的优势：速度最高、最大坏点较小、易部署到工程代码\n• XGBoost 在当前超参数下表现不佳，暂不作为最终方案\n\n组会建议表述：当前模型不是所有单项指标都第一，但在速度、尾部误差和部署便利性上更适合工程计算。", 23)

    s = blank_slide(prs, "可靠性门控：速度担心如何处理", "门控不等于每个点都调用 REFPROP")
    add_text(s, 0.75, 1.3, 6.1, 4.8, "轻量门控只做便宜检查：\n• 输入是否超出训练域\n• 正值物性是否为正\n• Cp 是否大于 Cv\n• 是否落入误差热力图高风险区\n\n只有失败点才 fallback 到 REFPROP。", 22)
    add_box(s, 7.2, 1.5, 4.4, 3.7, "ANN prediction\n↓\ncheap physics/range checks\n↓\npass: use ANN\nfail: call REFPROP", 24, (238, 242, 255), bold=True)
    add_text(s, 7.1, 5.55, 4.8, 0.7, "只要 fallback 比例很小，整体速度仍接近 ANN。", 21, (51, 65, 85))

    s = blank_slide(prs, "当前工作完成情况", "可作为论文主体结果")
    add_text(s, 0.75, 1.25, 11.6, 5.4, "已完成：\n• He4 单相 REFPROP 数据集整理\n• 多输出 ANN 模型训练与优化\n• log / shifted-log target 改进\n• 逐物性指标、R²、RMSE、P95/P99/max 统计\n• 温压平面误差热力图与 worst-point 分析\n• REFPROP 大规模速度对比\n• RF、Extra Trees、XGBoost、高斯核近似基线对比\n• 论文结果包整理", 23)

    s = blank_slide(prs, "下一步计划", "面向论文与工程应用")
    add_text(s, 0.75, 1.3, 11.5, 4.9, "短期：\n• 整理文献对比表\n• 增加物理一致性检查表\n• 进一步优化 MLP 或尝试小型 ensemble\n• 完成可靠性门控的速度影响评估\n\n中期：\n• 构建 5.5 K 以上单相氦4管路/换热器应用算例\n• 写论文初稿：数据、方法、误差、速度、基线对比", 23)

    return prs


def main():
    root = Path(__file__).resolve().parents[1]
    out = root / "paper_outputs" / "final_he4_single_phase" / "group_meeting_he4_ann_report_standard.pptx"
    prs = build(root)
    prs.save(out)
    shutil.copy2(out, root / "group_meeting_he4_ann_report_standard.pptx")
    print(f"Saved standard PPTX to: {out}")
    print(f"Copied to: {root / 'group_meeting_he4_ann_report_standard.pptx'}")


if __name__ == "__main__":
    main()
