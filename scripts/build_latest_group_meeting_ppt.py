from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PKG = ROOT / "paper_outputs" / "he4_2p2_hybrid_paper"
FIG = PKG / "figures"
TAB = PKG / "tables"


def rgb(hex_color: str) -> RGBColor:
    hex_color = hex_color.lstrip("#")
    return RGBColor(int(hex_color[0:2], 16), int(hex_color[2:4], 16), int(hex_color[4:6], 16))


def set_run(run, size=18, bold=False, color="#111827"):
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def add_textbox(slide, x, y, w, h, text, size=18, bold=False, color="#111827", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.space_after = Pt(4)
        for run in p.runs:
            set_run(run, size=size, bold=bold, color=color)
    return box


def add_bullets(slide, x, y, w, h, items, size=18, color="#1f2937"):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(8)
        p.font.name = "Microsoft YaHei"
        p.font.size = Pt(size)
        p.font.color.rgb = rgb(color)
    return box


def add_title(slide, title, subtitle=None):
    add_textbox(slide, 0.45, 0.18, 12.4, 0.46, title, size=24, bold=True, color="#111827")
    line = slide.shapes.add_shape(1, Inches(0.45), Inches(0.78), Inches(12.45), Inches(0.035))
    line.fill.solid()
    line.fill.fore_color.rgb = rgb("#2563eb")
    line.line.fill.background()
    if subtitle:
        add_textbox(slide, 0.48, 0.86, 12.2, 0.34, subtitle, size=12.5, color="#64748b")


def blank_slide(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_title(slide, title, subtitle)
    return slide


def add_card(slide, x, y, w, h, title, value, fill="#eff6ff", accent="#2563eb"):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb("#cbd5e1")
    bar = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.08), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(accent)
    bar.line.fill.background()
    add_textbox(slide, x + 0.18, y + 0.13, w - 0.28, 0.25, title, size=12.5, color="#475569")
    add_textbox(slide, x + 0.18, y + 0.44, w - 0.28, 0.46, value, size=23, bold=True, color="#111827")


def add_picture_fit(slide, path: Path, x, y, max_w, max_h):
    if not path.exists():
        add_textbox(slide, x, y, max_w, max_h, f"缺少图片:\n{path.name}", size=14, color="#b91c1c")
        return
    try:
        from PIL import Image

        with Image.open(path) as img:
            iw, ih = img.size
        scale = min(max_w / iw, max_h / ih)
        w = iw * scale
        h = ih * scale
        x2 = x + (max_w - w) / 2
        y2 = y + (max_h - h) / 2
        slide.shapes.add_picture(str(path), Inches(x2), Inches(y2), width=Inches(w), height=Inches(h))
    except Exception:
        slide.shapes.add_picture(str(path), Inches(x), Inches(y), width=Inches(max_w))


def add_simple_table(slide, df: pd.DataFrame, x, y, w, h, font_size=10.5):
    rows, cols = df.shape[0] + 1, df.shape[1]
    table = slide.shapes.add_table(rows, cols, Inches(x), Inches(y), Inches(w), Inches(h)).table
    for j, col in enumerate(df.columns):
        cell = table.cell(0, j)
        cell.text = str(col)
        cell.fill.solid()
        cell.fill.fore_color.rgb = rgb("#dbeafe")
        for p in cell.text_frame.paragraphs:
            for run in p.runs:
                set_run(run, size=font_size, bold=True)
    for i, (_, row) in enumerate(df.iterrows(), start=1):
        for j, col in enumerate(df.columns):
            val = row[col]
            if isinstance(val, float):
                text = f"{val:.4g}"
            else:
                text = str(val)
            cell = table.cell(i, j)
            cell.text = text
            for p in cell.text_frame.paragraphs:
                for run in p.runs:
                    set_run(run, size=font_size)
    return table


def build_ppt() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(13.333333)
    prs.slide_height = Inches(7.5)

    overall = pd.read_csv(TAB / "table_2_overall_ann_vs_hybrid.csv")
    per = pd.read_csv(TAB / "table_3_hybrid_per_property_metrics.csv")
    phase = pd.read_csv(TAB / "table_1_phase_counts.csv")
    speed = pd.read_csv(TAB / "table_5_speed_summary.csv")
    baseline = pd.read_csv(TAB / "table_4_baseline_summary.csv")
    ranges = pd.read_csv(TAB / "table_1_variable_ranges.csv")

    hybrid = overall[overall["model"] == "Local hybrid"].iloc[0]
    ann = overall[overall["model"] == "Critical-feature ANN"].iloc[0]
    speed_100k = speed[speed["n_points"] == 100000].iloc[0]

    # 1
    s = blank_slide(prs, "单相氦-4热物性神经网络代理模型最新进展", "组会汇报 | 2.2-300 K 单相 He-4 | Critical-feature ANN + Local Hybrid")
    add_textbox(s, 0.7, 1.4, 6.4, 0.9, "目标：建立一个比 REFPROP 更快、同时控制极端误差的氦-4单相物性快速预测模型", size=23, bold=True)
    add_card(s, 0.7, 3.0, 2.6, 1.05, "覆盖温度", "2.2-300 K", "#eff6ff", "#2563eb")
    add_card(s, 3.6, 3.0, 2.6, 1.05, "覆盖压力", "0.001-3 MPa", "#ecfdf5", "#16a34a")
    add_card(s, 6.5, 3.0, 2.6, 1.05, "输出物性", "8 个", "#fffbeb", "#d97706")
    add_card(s, 9.4, 3.0, 2.6, 1.05, "最终 Max", "4.65%", "#fef2f2", "#dc2626")
    add_bullets(s, 0.8, 5.0, 11.8, 1.3, [
        f"最终 hybrid：MAPE {hybrid['mape']:.4f}%，P99 {hybrid['p99']:.4f}%，within 5% = {hybrid['within_5pct']:.1f}%",
        f"100000 点批量测试：hybrid 比 REFPROP 快 {speed_100k['hybrid_speedup_vs_refprop']:.2f} 倍",
    ], size=19)

    # 2
    s = blank_slide(prs, "研究问题与路线", "为什么需要 ANN / hybrid 物性模型")
    add_bullets(s, 0.75, 1.25, 5.7, 5.5, [
        "REFPROP 精度高，但在大规模参数扫描和迭代计算中调用成本高",
        "低温氦物性在低温、近临界区、低密度区变化剧烈",
        "工程仿真不仅关心平均误差，更关心单个坏点导致的迭代风险",
        "本文路线：ANN 负责全域快速预测，local hybrid 负责压制高风险区域尾部误差",
    ], size=21)
    add_picture_fit(s, FIG / "fig_2_model_workflow.png", 6.7, 1.25, 5.9, 4.6)
    add_textbox(s, 6.9, 6.05, 5.5, 0.5, "Fig. 2 模型流程：REFPROP 数据 -> 单相筛选 -> 特征构造 -> ANN -> 局部修正 -> 八物性输出", size=13, color="#64748b")

    # 3
    s = blank_slide(prs, "数据集：2.2-300 K 单相氦-4", "REFPROP 生成 + 单相筛选")
    add_picture_fit(s, FIG / "fig_1_dataset_phase_map.png", 0.65, 1.2, 6.3, 5.45)
    add_card(s, 7.35, 1.35, 2.25, 0.95, "样本数", "358588", "#eff6ff", "#2563eb")
    add_card(s, 9.95, 1.35, 2.25, 0.95, "输入特征", "9 个", "#ecfdf5", "#16a34a")
    add_simple_table(s, phase[["phase_label", "n_samples", "fraction_pct"]].rename(columns={"phase_label": "Phase", "n_samples": "N", "fraction_pct": "%"}), 7.25, 2.75, 5.2, 2.05, font_size=10)
    add_bullets(s, 7.35, 5.25, 5.0, 1.1, [
        "去除两相区与近饱和风险点",
        "覆盖气相、液相、超临界等单相状态",
    ], size=17)

    # 4
    s = blank_slide(prs, "输入特征与输出物性", "加入相区标识和临界特征，降低低温/临界区建模难度")
    add_bullets(s, 0.75, 1.25, 5.9, 5.3, [
        "基础输入：Temperature, Pressure",
        "相区输入：phase_code",
        "临界特征：reduced_temperature, reduced_pressure, critical_distance, inverse_critical_distance 等",
        "目标变换：log / shifted-log + 标准化",
        "训练目标：多输出 MLP 同时预测八个物性",
    ], size=20)
    out_props = per[["property", "mape", "p99", "max"]].copy()
    out_props.columns = ["Property", "MAPE/%", "P99/%", "Max/%"]
    add_simple_table(s, out_props, 7.0, 1.25, 5.6, 5.35, font_size=9.4)

    # 5
    s = blank_slide(prs, "最终模型整体结果", "Hybrid 主要改善尾部误差")
    add_card(s, 0.75, 1.25, 2.35, 1.05, "Hybrid MAPE", f"{hybrid['mape']:.4f}%", "#eff6ff", "#2563eb")
    add_card(s, 3.35, 1.25, 2.35, 1.05, "Hybrid P99", f"{hybrid['p99']:.4f}%", "#fffbeb", "#d97706")
    add_card(s, 5.95, 1.25, 2.35, 1.05, "Hybrid Max", f"{hybrid['max']:.4f}%", "#fef2f2", "#dc2626")
    add_card(s, 8.55, 1.25, 2.35, 1.05, "Mean R2", f"{hybrid['mean_r2']:.8f}", "#ecfdf5", "#16a34a")
    compare = overall[["model", "mape", "p99", "max", "within_5pct"]].copy()
    compare.columns = ["Model", "MAPE/%", "P99/%", "Max/%", "Within 5%"]
    add_simple_table(s, compare, 0.9, 3.0, 5.7, 1.35, font_size=10.5)
    add_picture_fit(s, FIG / "fig_5_ann_vs_hybrid_overall_errors.png", 6.9, 2.65, 5.4, 3.1)
    add_bullets(s, 0.95, 5.3, 11.3, 1.1, [
        f"纯 ANN 最大误差 {ann['max']:.2f}%，hybrid 后降至 {hybrid['max']:.2f}%",
        "所有测试点误差均小于 5%，直接回应工程计算中的坏点风险",
    ], size=18)

    # 6
    s = blank_slide(prs, "逐物性精度", "所有输出物性最大误差均小于 5%")
    add_picture_fit(s, FIG / "fig_3_per_property_error_metrics.png", 0.55, 1.2, 6.35, 5.55)
    add_picture_fit(s, FIG / "fig_3_r2_and_within_thresholds.png", 6.95, 1.2, 5.95, 5.55)

    # 7
    s = blank_slide(prs, "误差在温压平面上的分布", "残余误差主要集中在低温与近临界区域")
    add_picture_fit(s, FIG / "fig_4_hybrid_error_maps" / "fig_4a_hybrid_max_error_heatmap.png", 0.65, 1.15, 5.8, 4.85)
    add_picture_fit(s, FIG / "fig_4_hybrid_error_maps" / "fig_4b_hybrid_mean_error_heatmap.png", 6.85, 1.15, 5.8, 4.85)
    add_textbox(s, 0.85, 6.25, 11.7, 0.55, "结论：hybrid 后整体误差较低，剩余高误差区域与氦物性本身的强非线性区域一致。", size=18, color="#334155")

    # 8
    s = blank_slide(prs, "Hybrid 修正效果", "用小幅速度代价压制尾部坏点")
    add_picture_fit(s, FIG / "fig_5_ann_vs_hybrid_max_error.png", 0.65, 1.15, 6.0, 4.95)
    add_picture_fit(s, FIG / "fig_5_hybrid_worst_points_tp.png", 7.0, 1.15, 5.4, 4.95)
    add_textbox(s, 0.9, 6.25, 11.4, 0.52, "核心表述：hybrid 不是每点调用 REFPROP，而是在预定义高风险区域进行局部插值修正。", size=18, color="#334155")

    # 9
    s = blank_slide(prs, "与传统机器学习模型对比", "不能只看平均 MAPE，需要同时看 Max 和速度")
    add_picture_fit(s, FIG / "fig_6_baseline_accuracy_speed_tail.png", 0.55, 1.1, 6.4, 5.55)
    bshow = baseline[["model", "mape", "max", "points_per_s"]].copy()
    bshow.columns = ["Model", "MAPE/%", "Max/%", "points/s"]
    add_simple_table(s, bshow, 7.1, 1.35, 5.4, 3.4, font_size=9.0)
    add_bullets(s, 7.2, 5.15, 5.1, 1.1, [
        "树模型平均误差可更低，但最大误差更大且速度更慢",
        "最终 hybrid 的优势是尾部可靠性 + 速度综合平衡",
    ], size=16.5)

    # 10
    s = blank_slide(prs, "与 REFPROP 速度对比", "100000 点测试中 hybrid 仍有明显速度优势")
    add_picture_fit(s, FIG / "fig_7_speed_benchmark.png", 0.65, 1.15, 6.35, 5.35)
    add_card(s, 7.35, 1.35, 2.25, 0.95, "ANN", f"{speed_100k['ann_points_per_s']:.0f}\npoints/s", "#eff6ff", "#2563eb")
    add_card(s, 9.95, 1.35, 2.25, 0.95, "Hybrid", f"{speed_100k['hybrid_points_per_s']:.0f}\npoints/s", "#ecfdf5", "#16a34a")
    add_card(s, 7.35, 3.0, 2.25, 0.95, "REFPROP", f"{speed_100k['refprop_points_per_s']:.0f}\npoints/s", "#fef2f2", "#dc2626")
    add_card(s, 9.95, 3.0, 2.25, 0.95, "Hybrid 加速", f"{speed_100k['hybrid_speedup_vs_refprop']:.2f}x", "#fffbeb", "#d97706")
    add_textbox(s, 7.35, 4.85, 4.9, 1.05, f"Hybrid 相比纯 ANN 只慢 {(speed_100k['hybrid_slowdown_vs_ann'] - 1) * 100:.2f}%，但把最大误差压到 5% 以下。", size=18, bold=True, color="#111827")

    # 11
    s = blank_slide(prs, "预生成数据库插值的补充思考", "插值很快，但尾部可靠性不足")
    add_bullets(s, 0.75, 1.25, 6.0, 4.8, [
        "用 22 万预生成点 + KDTree 插值做过初步测试",
        "插值速度可达 ANN 的数倍甚至十几倍",
        "但最大误差显著偏大，部分设置超过 100%",
        "原因：低温、近临界、近相界区域物性变化太剧烈，普通插值很难全域保证尾部误差",
    ], size=20)
    interp = pd.DataFrame(
        [
            ["ANN", "2.55e5", "0.165", "46.1"],
            ["KDTree k=8", "1.18e6", "0.164", "90.6"],
            ["Hybrid", "2.37e5", "0.131", "4.65"],
        ],
        columns=["Method", "points/s", "MAPE/%", "Max/%"],
    )
    add_simple_table(s, interp, 7.2, 1.65, 4.75, 1.8, font_size=13)
    add_textbox(s, 7.2, 4.25, 4.9, 1.25, "可作为论文补充基线：说明本文方法不是绝对最快，而是速度与可靠性更均衡。", size=19, bold=True, color="#334155")

    # 12
    s = blank_slide(prs, "论文主线与下一步", "建议收束为高可靠单相 He-4 物性代理模型")
    add_bullets(s, 0.75, 1.2, 6.0, 5.6, [
        "论文主旨：面向低温工程计算的高可靠 He-4 单相物性代理模型",
        "核心卖点：2.2-300 K、8 个物性、所有物性最大误差 < 5%",
        "不建议把 GM/稀释制冷机算例作为正文核心，避免转移审稿焦点",
        "可补强：正式英文摘要、方法流程图润色、参考文献对比、图表统一格式",
    ], size=20)
    add_card(s, 7.35, 1.3, 4.0, 1.0, "一句话结论", "17.45x faster than REFPROP", "#eff6ff", "#2563eb")
    add_card(s, 7.35, 2.85, 4.0, 1.0, "可靠性", "Max error < 5%", "#ecfdf5", "#16a34a")
    add_card(s, 7.35, 4.4, 4.0, 1.0, "论文定位", "fast and reliable surrogate", "#fffbeb", "#d97706")

    return prs


def main() -> None:
    prs = build_ppt()
    out = PKG / "group_meeting_latest_he4_ann_hybrid.pptx"
    prs.save(out)
    root_copy = ROOT / "group_meeting_latest_he4_ann_hybrid.pptx"
    prs.save(root_copy)
    print(f"Saved: {out}")
    print(f"Saved: {root_copy}")


if __name__ == "__main__":
    main()
