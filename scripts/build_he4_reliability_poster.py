from __future__ import annotations

from pathlib import Path

import pandas as pd
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
PKG = ROOT / "paper_outputs" / "he4_2p2_hybrid_paper"
FIG = PKG / "figures"
TAB = PKG / "tables"


def rgb(value: str) -> RGBColor:
    value = value.lstrip("#")
    return RGBColor(int(value[:2], 16), int(value[2:4], 16), int(value[4:], 16))


def set_font(run, size=14, bold=False, color="#111827"):
    run.font.name = "Microsoft YaHei"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = rgb(color)


def text_box(slide, x, y, w, h, text, size=14, bold=False, color="#111827", align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.TOP
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.alignment = align
        p.space_after = Pt(2)
        for run in p.runs:
            set_font(run, size=size, bold=bold, color=color)
    return box


def panel(slide, x, y, w, h, title, fill="#ffffff", line="#d1d5db"):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb(line)
    text_box(slide, x + 0.18, y + 0.12, w - 0.36, 0.3, title, size=14, bold=True, color="#0f172a")
    return shape


def card(slide, x, y, w, h, title, value, fill, accent):
    shape = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.color.rgb = rgb("#cbd5e1")
    bar = slide.shapes.add_shape(1, Inches(x), Inches(y), Inches(0.07), Inches(h))
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(accent)
    bar.line.fill.background()
    text_box(slide, x + 0.14, y + 0.10, w - 0.22, 0.22, title, size=8.8, color="#475569")
    text_box(slide, x + 0.14, y + 0.34, w - 0.22, 0.34, value, size=16, bold=True, color="#111827")


def picture(slide, path: Path, x, y, max_w, max_h):
    if not path.exists():
        text_box(slide, x, y, max_w, max_h, f"缺少图片: {path.name}", size=10, color="#b91c1c")
        return
    from PIL import Image

    with Image.open(path) as img:
        iw, ih = img.size
    scale = min(max_w / iw, max_h / ih)
    w = iw * scale
    h = ih * scale
    slide.shapes.add_picture(
        str(path),
        Inches(x + (max_w - w) / 2),
        Inches(y + (max_h - h) / 2),
        width=Inches(w),
        height=Inches(h),
    )


def build() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(16)
    prs.slide_height = Inches(9)
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    # Background
    bg = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(16), Inches(9))
    bg.fill.solid()
    bg.fill.fore_color.rgb = rgb("#f8fafc")
    bg.line.fill.background()

    overall = pd.read_csv(TAB / "table_2_overall_ann_vs_hybrid.csv")
    hybrid = overall[overall["model"] == "Local hybrid"].iloc[0]
    speed = pd.read_csv(TAB / "table_5_speed_summary.csv")
    speed_100k = speed[speed["n_points"] == 100000].iloc[0]

    # Header
    text_box(
        slide,
        0.45,
        0.22,
        11.8,
        0.42,
        "低温氦物性机器学习代理模型的可靠性与误差控制",
        size=25,
        bold=True,
        color="#0f172a",
    )
    text_box(
        slide,
        0.48,
        0.72,
        11.4,
        0.28,
        "从“替代 REFPROP 加速”重新定位为：以 REFPROP 为基准，研究低温强非线性物性代理模型的坏点、尾部误差和可靠性",
        size=11.5,
        color="#475569",
    )
    card(slide, 12.25, 0.18, 1.05, 0.72, "MAPE", f"{hybrid['mape']:.3f}%", "#eff6ff", "#2563eb")
    card(slide, 13.42, 0.18, 1.05, 0.72, "Max", f"{hybrid['max']:.2f}%", "#fef2f2", "#dc2626")
    card(slide, 14.59, 0.18, 1.05, 0.72, "Speedup", f"{speed_100k['hybrid_speedup_vs_refprop']:.1f}x", "#ecfdf5", "#16a34a")

    # Left column
    panel(slide, 0.45, 1.18, 3.45, 2.05, "研究问题：老师质疑后的新定位", "#ffffff")
    text_box(
        slide,
        0.65,
        1.62,
        3.05,
        1.35,
        "不再声称：ANN 替代 REFPROP。\n\n本文更合理的问题：\n1. 低温/临界区 ML 物性模型哪里会失效？\n2. 平均误差很低时，坏点是否仍危险？\n3. 如何控制 P99 和最大误差？",
        size=10.8,
        color="#1f2937",
    )

    panel(slide, 0.45, 3.43, 3.45, 2.12, "数据集", "#ffffff")
    text_box(
        slide,
        0.65,
        3.86,
        3.05,
        1.42,
        "工质：单相 helium-4\n数据源：NIST REFPROP\n温度：2.2-300 K\n压力：0.001-3 MPa\n样本数：358588\n输出：8 个热物性/输运物性",
        size=11.2,
        color="#1f2937",
    )

    panel(slide, 0.45, 5.74, 3.45, 2.73, "讲给老师听的一句话", "#eef2ff", "#bfdbfe")
    text_box(
        slide,
        0.68,
        6.18,
        3.0,
        1.85,
        "REFPROP 是权威基准。\n我的工作不是证明 ANN 更“正确”，而是研究：\n\n在低温复杂区域，机器学习代理模型如何识别和压制会破坏工程迭代的极端坏点。",
        size=12.0,
        bold=True,
        color="#1e3a8a",
    )

    # Center column
    panel(slide, 4.18, 1.18, 5.55, 2.55, "方法流程", "#ffffff")
    picture(slide, FIG / "fig_2_model_workflow.png", 4.35, 1.62, 5.18, 1.78)

    panel(slide, 4.18, 3.98, 2.68, 2.22, "数据域与相区", "#ffffff")
    picture(slide, FIG / "fig_1_dataset_phase_map.png", 4.33, 4.39, 2.35, 1.53)

    panel(slide, 7.05, 3.98, 2.68, 2.22, "尾部误差修正", "#ffffff")
    picture(slide, FIG / "fig_5_ann_vs_hybrid_max_error.png", 7.18, 4.39, 2.42, 1.53)

    panel(slide, 4.18, 6.48, 5.55, 1.99, "关键方法", "#ffffff")
    text_box(
        slide,
        4.42,
        6.88,
        5.05,
        1.18,
        "Critical-feature ANN：T, P + phase_code + reduced/critical features。\nLocal hybrid correction：只在低温液体、近临界、低密度等高风险区做局部修正。\n目标：不是最低平均 MAPE，而是让所有物性 Max error < 5%。",
        size=10.7,
        color="#1f2937",
    )

    # Right column
    panel(slide, 10.02, 1.18, 5.53, 2.28, "核心结果", "#ffffff")
    card(slide, 10.28, 1.66, 1.22, 0.82, "P99", f"{hybrid['p99']:.3f}%", "#fffbeb", "#d97706")
    card(slide, 11.72, 1.66, 1.22, 0.82, "Within 5%", f"{hybrid['within_5pct']:.0f}%", "#ecfdf5", "#16a34a")
    card(slide, 13.16, 1.66, 1.22, 0.82, "R2", f"{hybrid['mean_r2']:.6f}", "#eff6ff", "#2563eb")
    text_box(
        slide,
        10.28,
        2.64,
        4.95,
        0.45,
        "Hybrid 将纯 ANN 最大误差从 46.10% 压到 4.65%，所有输出物性的最大误差均小于 5%。",
        size=11.2,
        bold=True,
        color="#111827",
    )

    panel(slide, 10.02, 3.72, 2.68, 2.13, "逐物性误差", "#ffffff")
    picture(slide, FIG / "fig_3_per_property_error_metrics.png", 10.17, 4.15, 2.38, 1.42)

    panel(slide, 12.88, 3.72, 2.67, 2.13, "速度对比", "#ffffff")
    picture(slide, FIG / "fig_7_speed_benchmark.png", 13.02, 4.15, 2.38, 1.42)

    panel(slide, 10.02, 6.12, 5.53, 2.35, "下一步：让工作更有意义", "#fff7ed", "#fed7aa")
    text_box(
        slide,
        10.26,
        6.55,
        5.05,
        1.42,
        "1. 弱化“加速 REFPROP”，强化可靠性/坏点/不确定性。\n2. 增加 uncertainty-aware 实验：ensemble 或 MC dropout 识别高误差区域。\n3. 小规模多工质 proof-of-concept：He, H2, N2, Ar, Ne，用对应态参数探索泛化。",
        size=10.8,
        bold=True,
        color="#7c2d12",
    )

    # Footer
    text_box(
        slide,
        0.48,
        8.62,
        15.1,
        0.22,
        "当前结论：He-4 结果可作为 case study 保留，但论文主旨应转为“低温物性代理模型可靠性研究”，后续以不确定性和多工质泛化提升意义。",
        size=9.6,
        color="#64748b",
        align=PP_ALIGN.CENTER,
    )

    return prs


def main() -> None:
    prs = build()
    out = PKG / "he4_reliability_poster_for_group_meeting.pptx"
    root_out = ROOT / "he4_reliability_poster_for_group_meeting.pptx"
    prs.save(out)
    prs.save(root_out)
    print(f"Saved: {out}")
    print(f"Saved: {root_out}")


if __name__ == "__main__":
    main()
